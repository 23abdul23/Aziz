# ...existing code...
# jarvis/extract.py
from pathlib import Path
import os

# Text formats
def read_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")

# PDF
import fitz  # PyMuPDF
def read_pdf(path: Path) -> str:
    doc = fitz.open(path)
    text = []
    for page in doc:
        text.append(page.get_text())
    return "\n".join(text)

# DOCX
import docx
def read_docx(path: Path) -> str:
    doc = docx.Document(path)
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paras)

# PPTX
from pptx import Presentation
def read_pptx(path: Path) -> str:
    prs = Presentation(path)
    slides_text = []
    for slide in prs.slides:
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_content.append(shape.text)
        slides_text.append("\n".join(slide_content))
    return "\n\n".join(slides_text)

# Markdown (read as plain text)
def read_md(path: Path) -> str:
    return read_text_file(path)

# Main extractor
READERS = {
    ".txt": read_text_file,
    ".md": read_md,
    ".pdf": read_pdf,
    ".docx": read_docx,
    ".pptx": read_pptx
}

def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in READERS:
        try:
            return READERS[ext](path)
        except Exception as e:
            print(f"[WARN] Failed to read {path}: {e}")
            return ""
    else:
        return ""
